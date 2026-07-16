#!/usr/bin/env python3
"""Build RiyalAI presentation from the official امد 2026 student template."""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

BASE = Path("/Users/jenanhamad/Downloads/عروض الطلاب | امد 2026.pptx")
OUTPUT = Path("/Users/jenanhamad/Downloads/RiyalAI_امد2026.pptx")
SCREENSHOTS_DIR = Path(__file__).resolve().parent.parent / "presentation" / "screenshots"

# Slide 12 screenshot grid (2 rows: 4 then 3), positions borrowed from a
# previously hand-laid-out deck since the blank template has no photo slots there.
SCREENSHOT_SLOTS = [
    ("01-voice.png", Emu(537376), Emu(841248), Emu(1614280), Emu(1993392)),
    ("02-voice-recording.png", Emu(2689032), Emu(841248), Emu(1614280), Emu(1993392)),
    ("03-receipt.png", Emu(4840688), Emu(841248), Emu(1614280), Emu(1993392)),
    ("04-add-expense.png", Emu(6992344), Emu(841248), Emu(1614280), Emu(1993392)),
    ("05-home.png", Emu(1075290), Emu(2944368), Emu(1614280), Emu(1993392)),
    ("06-friends.png", Emu(3764860), Emu(2944368), Emu(1614280), Emu(1993392)),
    ("07-leaderboard.png", Emu(6454430), Emu(2944368), Emu(1614280), Emu(1993392)),
]


def set_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = text
        return
    tf.paragraphs[0].text = text
    for p in tf.paragraphs[1:]:
        p.text = ""


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)


def find_shape_by_name(slide, name: str):
    for shape in walk_shapes(slide.shapes):
        if shape.name == name:
            return shape
    return None


def set_named(slide, name: str, text: str) -> None:
    shape = find_shape_by_name(slide, name)
    if shape:
        set_text(shape, text)


def add_slide_screenshots(slide) -> None:
    """Lay app screenshots out in a 4+3 grid on the demo slide."""
    for filename, left, top, width, height in SCREENSHOT_SLOTS:
        img_path = SCREENSHOTS_DIR / filename
        if not img_path.exists():
            continue
        slide.shapes.add_picture(str(img_path), left, top, width, height)


def remove_shape(slide, name: str) -> None:
    shape = find_shape_by_name(slide, name)
    if shape is not None:
        shape._element.getparent().remove(shape._element)


CONTENT = {
    # Slide 1
    "s1_title": "ريـال AI",
    "s1_subtitle": "ريالي — تتبع مصاريف بالذكاء الاصطناعي",
    # Slide 2
    "s2_m1": "جنان الغطيمل",
    "s2_m2": "الهنوف الجمعان",
    "s2_m3": "سارة الزهراني",
    # Slide 4 — README.md + App.js
    "s4_title": "المشكلة  وحلّها",
    "s4_desc": (
        "كثير من الناس ما يتابعون مصاريفهم يومياً، وصاحبات المشاريع الصغيرة يحتاجون "
        "يربطون الإيراد والمصروف ويشوفون الربح والضريبة — بدون تعقيد. ريالي يجمع "
        "الاثنين في تطبيق واحد: وضع أفراد (تحفيز وتلعيب) ووضع أعمال (ربح وضريبة وهدر)."
    ),
    "s4_box1": "تسجيل بالصوت أو صورة الإيصال — AI يستخرج المبلغ والفئة تلقائياً",
    "s4_box2": "وضعان في نفس التطبيق: أفراد (XP وتحديات) أو أعمال (لوحة ربح)",
    "s4_box3": "واجهة عربية RTL، PWA يعمل على الجوال والويب",
    "s4_lbl1": "إدخال ذكي",
    "s4_lbl2": "وضعان",
    "s4_lbl3": "عربي RTL",
    # Slide 5 — database.py
    "s5_title": "البيانات المستخدمة",
    "s5_desc": (
        "البيانات مخزّنة في SQLite: مصاريف (مبلغ، فئة، تاريخ، mode، entry_type)، "
        "ملفات مستخدم (XP، مستوى، سلسلة، active_mode)، تحديات، صداقات، وتوكنات استعادة كلمة المرور."
    ),
    "s5_bullets": (
        "مصادر الإدخال: يدوي، صوت عربي (/voice/process)، صورة إيصال (/receipt/process)\n"
        "المعالجة: OpenRouter — Claude للتحديات والنص، Gemini للصوت والإيصالات\n"
        "التحديات التقنية: دقة اللهجة السعودية، وفصل بيانات personal عن business عبر حقل mode"
    ),
    # Slide 6 — Dockerfile + openrouter.py + auth.py
    "tech": [
        ("Google Shape;144;p18", "React (PWA)\nواجهة RTL — Home، Voice، Challenges، Friends"),
        ("Google Shape;149;p18", "FastAPI + SQLite\nAPI موحّد يخدم الواجهة ويخزّن البيانات"),
        ("Google Shape;155;p18", "OpenRouter\nClaude Haiku — تحديات أسبوعية بالعربي"),
        ("Google Shape;160;p18", "Gemini Vision\nOCR إيصالات + تفريغ الصوت العربي"),
        ("Google Shape;166;p18", "Docker → Railway\nحاوية واحدة: API + React build"),
        ("Google Shape;171;p18", "JWT + SMTP\nتسجيل دخول واستعادة كلمة المرور بالبريد"),
    ],
    # Slide 7 — App.js + المكوّنات
    "s7_title": "ريالي — تطبيق واحد، وضعان",
    "s7_subtitle": "وصف الفكرة",
    "s7_desc": (
        "PWA عربي لتتبع المصاريف بالذكاء الاصطناعي. وضع الأفراد: XP (+20 لكل مصروف)، "
        "سلسلة يومية (×2 باليوم 7+)، تحديات AI أسبوعية، متصدرين، أصدقاء، قصة أسبوعك، وتحليلات. "
        "وضع الأعمال: لوحة ربح، صحة المشروع، تقدير VAT، كاشف هدر، ونظرة على المشروع. "
        "التبديل بين الوضعين من ModeSwitcher في أي وقت."
    ),
    # Slide 8 — مسارات API من main.py
    "s8_title": "كيفية توفير هذه البيانات وكيفية استخدامها",
    "s8_flow": (
        "١. المستخدم يسجّل حساباً ويختار وضع أفراد أو أعمال عند التسجيل\n"
        "٢. يُدخل مصروفه: صوت (/voice) أو إيصال (/receipt) أو يدوي (/add)\n"
        "٣. OpenRouter يستخرج المبلغ والفئة → المستخدم يؤكد → يُحفظ في expenses\n"
        "٤. وضع personal: XP + تحديات + leaderboard | وضع business: dashboard + glance\n"
        "٥. البيانات تُعرض في لوحات React حسب active_mode المخزّن في users"
    ),
    # Slide 9
    "s9_title": "مواءمة الفكرة :",
    "s9_desc": (
        "ريالي يستخدم الذكاء الاصطناعي لتحسين الوعي المالي:\n"
        "• توليد تحديات أسبوعية مخصصة من أنماط الإنفاق (OpenRouter/Claude)\n"
        "• فهم الصوت العربي المحكي وتحويله لمصروف جاهز (Gemini)\n"
        "• قراءة الإيصالات بالصورة واستخراج البيانات (Gemini Vision)\n"
        "• تحليل أعمال: ربح، تقدير ضريبة، وكشف هدر — من business.py"
    ),
    # Slide 10
    "s10_title": "ملخص",
    "s10_desc": (
        "ريالي تطبيق منشور على riyalai.up.railway.app — React PWA + FastAPI + SQLite "
        "في حاوية Docker واحدة. يدعم إدخال المصاريف بالصوت والإيصال واليدوي، مع وضعين: "
        "أفراد (تلعيب وتحديات ومتصدرين) وأعمال (ربح وVAT وهدر). مصادقة JWT، استعادة "
        "كلمة المرور عبر SMTP، وواجهة عربية كاملة."
    ),
    # Slide 11
    "s11_title": "الاختبار/التحقق:",
    "s11_desc": (
        "✅ منشور على: riyalai.up.railway.app (Railway + Docker)\n"
        "✅ تسجيل صوتي + OCR إيصالات + إدخال يدوي — يعمل عبر API\n"
        "✅ وضع أعمال: /business/dashboard و /business/glance\n"
        "✅ تلعيب: XP، مستويات، تحديات AI، leaderboard، أصدقاء\n"
        "✅ نسيت كلمة المرور: /auth/forgot-password + /reset-password"
    ),
    # Slide 12
    "s12_title": "العرض التوضيحي",
    "s12_desc": "https://riyalai.up.railway.app  —  PWA: iOS & Android & Web  —  جاهز للعرض ✅",
    # Slide 13 — ما هو مبني vs ما يمكن تطويره
    "s13_help": (
        "ما تحتاج إلى مساعدة فيه:\n"
        "نماذج صوتية أدق للهجة السعودية، ومراجعة قانونية لتقدير VAT."
    ),
    "s13_challenges": (
        "التحديات:\n"
        "• دقة التعرف على اللهجات في التسجيل الصوتي\n"
        "• تقدير VAT تقديري وليس استشارة ضريبية رسمية\n"
        "• كاشف الهدر يعتمد قواعد + AI حسب توفر المفتاح"
    ),
    "s13_future": (
        "الخطة المستقبلية:\n"
        "• إشعارات عند تجاوز الميزانية\n"
        "• تحليلات أعمق ورسوم بيانية تفاعلية\n"
        "• دعم حسابات أعمال متعددة لكل مستخدم"
    ),
}


def build() -> Path:
    shutil.copy2(BASE, OUTPUT)
    prs = Presentation(str(OUTPUT))

    # Slide 1
    s = prs.slides[0]
    set_named(s, "Google Shape;50;p13", CONTENT["s1_title"])
    set_named(s, "Google Shape;51;p13", CONTENT["s1_subtitle"])

    # Slide 2 — template has 4 member slots, we only have 3: drop the 4th
    # (leftmost in RTL reading order: name box 65, icon circle 66, icon pic 68)
    s = prs.slides[1]
    set_named(s, "Google Shape;62;p14", CONTENT["s2_m1"])
    set_named(s, "Google Shape;63;p14", CONTENT["s2_m2"])
    set_named(s, "Google Shape;64;p14", CONTENT["s2_m3"])
    remove_shape(s, "Google Shape;65;p14")
    remove_shape(s, "Google Shape;66;p14")
    remove_shape(s, "Google Shape;68;p14")

    # Slide 4
    s = prs.slides[3]
    set_named(s, "Google Shape;99;p16", CONTENT["s4_desc"])
    set_named(s, "Google Shape;100;p16", CONTENT["s4_box1"])
    set_named(s, "Google Shape;101;p16", CONTENT["s4_box2"])
    set_named(s, "Google Shape;102;p16", CONTENT["s4_box3"])
    set_named(s, "Google Shape;103;p16", CONTENT["s4_lbl1"])
    set_named(s, "Google Shape;104;p16", CONTENT["s4_lbl2"])
    set_named(s, "Google Shape;105;p16", CONTENT["s4_lbl3"])

    # Slide 5
    s = prs.slides[4]
    set_named(s, "Google Shape;121;p17", CONTENT["s5_desc"])
    set_named(s, "Google Shape;122;p17", CONTENT["s5_bullets"])

    # Slide 6
    s = prs.slides[5]
    for shape_name, text in CONTENT["tech"]:
        set_named(s, shape_name, text)

    # Slide 7
    s = prs.slides[6]
    set_named(s, "Google Shape;182;p19", CONTENT["s7_desc"])
    set_named(s, "Google Shape;183;p19", CONTENT["s7_title"])
    set_named(s, "Google Shape;184;p19", CONTENT["s7_subtitle"])

    # Slide 8
    s = prs.slides[7]
    set_named(s, "Google Shape;190;p20", CONTENT["s8_flow"])
    set_named(s, "Google Shape;191;p20", CONTENT["s8_title"])

    # Slide 9
    s = prs.slides[8]
    set_named(s, "Google Shape;200;p21", CONTENT["s9_desc"])

    # Slide 10
    s = prs.slides[9]
    set_named(s, "Google Shape;209;p22", CONTENT["s10_desc"])

    # Slide 11
    s = prs.slides[10]
    set_named(s, "Google Shape;216;p23", CONTENT["s11_desc"])

    # Slide 12 — demo + screenshots. Template only ships a title + one
    # instructions paragraph here, so shrink/reposition the title to make
    # room and drop screenshots into a 4+3 grid beneath it.
    s = prs.slides[11]
    title_shape = find_shape_by_name(s, "Google Shape;222;p24")
    if title_shape is not None:
        set_text(title_shape, CONTENT["s12_title"])
        title_shape.left = Emu(2323500)
        title_shape.top = Emu(-9144)
        title_shape.width = Emu(6363300)
        title_shape.height = Emu(502920)
    desc_shape = find_shape_by_name(s, "Google Shape;223;p24")
    if desc_shape is not None:
        set_text(desc_shape, CONTENT["s12_desc"])
        desc_shape.left = Emu(457200)
        desc_shape.top = Emu(530352)
        desc_shape.width = Emu(8229600)
        desc_shape.height = Emu(365760)
    add_slide_screenshots(s)

    # Slide 13 — three top-level groups (help, challenges, future)
    s = prs.slides[12]
    group_map = {
        "Google Shape;230;p25": CONTENT["s13_help"],
        "Google Shape;233;p25": CONTENT["s13_challenges"],
        "Google Shape;236;p25": CONTENT["s13_future"],
    }
    for shape in s.shapes:
        if shape.name in group_map and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                if sub.has_text_frame:
                    set_text(sub, group_map[shape.name])
                    break

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Saved: {out}")
