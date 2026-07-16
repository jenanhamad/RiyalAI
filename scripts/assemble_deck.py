#!/usr/bin/env python3
"""Assemble rendered slide PNGs into a 16:9 PPTX and a PDF."""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

SLIDES_DIR = Path(__file__).resolve().parent.parent / "presentation" / "deck" / "slides"
PPTX_OUT = Path.home() / "Downloads" / "RiyalAI_عرض_لجنة_التحكيم.pptx"
PDF_OUT = Path.home() / "Downloads" / "RiyalAI_عرض_لجنة_التحكيم.pdf"


def build() -> None:
    pngs = sorted(SLIDES_DIR.glob("slide*.png"))
    if not pngs:
        raise SystemExit("No rendered slides found — run scripts/render_deck.mjs first")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(PPTX_OUT))
    print(f"Saved: {PPTX_OUT}")

    images = [Image.open(p).convert("RGB") for p in pngs]
    images[0].save(str(PDF_OUT), save_all=True, append_images=images[1:], resolution=144)
    print(f"Saved: {PDF_OUT}")


if __name__ == "__main__":
    build()
